#!/usr/bin/env python3
"""
TBWA Creative Campaign Analysis System - PageIndex Agent
ColPali-style file processing with Azure OpenAI integration

This agent processes creative files (PDFs, images, videos, presentations) 
and creates semantic chunks with embeddings for searchable indexing.
"""

import os
import sys
import json
import asyncio
import logging
import argparse
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
from dataclasses import dataclass, asdict
from datetime import datetime
import uuid

# Core dependencies
import numpy as np
import pandas as pd
from PIL import Image
import cv2
import requests
from tqdm import tqdm

# Azure AI
from openai import AzureOpenAI
from azure.identity import DefaultAzureCredential, ClientSecretCredential
from azure.storage.blob import BlobServiceClient

# Document processing
import fitz  # PyMuPDF for PDF processing
from docx import Document
from pptx import Presentation

# Environment
from dotenv import load_dotenv
import pyodbc

# Load environment variables
load_dotenv()

@dataclass
class ChunkMetadata:
    """Metadata for a semantic chunk"""
    chunk_id: str
    file_id: str
    filepath: str
    title: Optional[str]
    text_snippet: str
    tags: List[str]
    visual_quality_score: float
    semantic_topics: List[str]
    mood_label: str
    chunk_index: int
    chunk_size: int
    content_type: str  # 'text', 'image', 'slide', 'video'
    confidence_score: float
    embedding_vector: Optional[List[float]] = None

@dataclass
class FileMetadata:
    """Metadata for a processed file"""
    file_id: str
    original_filename: str
    filepath: str
    file_size: int
    mime_type: str
    campaign_name: Optional[str]
    client_name: Optional[str]
    brand_name: Optional[str]
    file_type: str  # 'video', 'image', 'presentation', 'document'
    processing_status: str = 'pending'
    total_chunks: int = 0
    azure_blob_url: Optional[str] = None
    google_drive_id: Optional[str] = None
    error_message: Optional[str] = None

class PageIndexAgent:
    """
    ColPali-style PageIndex agent for semantic file processing
    """
    
    def __init__(self):
        self.setup_logging()
        self.setup_azure_clients()
        self.setup_database()
        
    def setup_logging(self):
        """Configure logging"""
        logging.basicConfig(
            level=logging.INFO,
            format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
            handlers=[
                logging.FileHandler('pageindex_agent.log'),
                logging.StreamHandler(sys.stdout)
            ]
        )
        self.logger = logging.getLogger(__name__)
        
    def setup_azure_clients(self):
        """Initialize Azure clients"""
        try:
            # Azure OpenAI
            self.openai_client = AzureOpenAI(
                api_key=os.getenv('AZURE_OPENAI_API_KEY'),
                api_version=os.getenv('AZURE_OPENAI_API_VERSION', '2024-02-15-preview'),
                azure_endpoint=os.getenv('AZURE_OPENAI_ENDPOINT')
            )
            
            # Azure Blob Storage
            if os.getenv('AZURE_STORAGE_CONNECTION_STRING'):
                self.blob_client = BlobServiceClient.from_connection_string(
                    os.getenv('AZURE_STORAGE_CONNECTION_STRING')
                )
            
            self.logger.info("✅ Azure clients initialized successfully")
            
        except Exception as e:
            self.logger.error(f"❌ Failed to initialize Azure clients: {e}")
            raise
            
    def setup_database(self):
        """Initialize database connection"""
        try:
            # SQL Server connection string
            server = os.getenv('AZURE_SQL_SERVER', 'sqltbwaprojectscoutserver.database.windows.net')
            database = os.getenv('AZURE_SQL_DATABASE', 'SQL-TBWA-ProjectScout-Reporting-Prod')
            username = os.getenv('AZURE_SQL_USER', 'sqladmin')
            password = os.getenv('AZURE_SQL_PASSWORD')
            
            if not password:
                raise ValueError("AZURE_SQL_PASSWORD environment variable is required")
                
            self.connection_string = (
                f"DRIVER={{ODBC Driver 18 for SQL Server}};"
                f"SERVER={server};"
                f"DATABASE={database};"
                f"UID={username};"
                f"PWD={password};"
                f"Encrypt=yes;"
                f"TrustServerCertificate=no;"
                f"Connection Timeout=30;"
            )
            
            # Test connection
            with pyodbc.connect(self.connection_string) as conn:
                cursor = conn.cursor()
                cursor.execute("SELECT 1")
                
            self.logger.info("✅ Database connection established")
            
        except Exception as e:
            self.logger.error(f"❌ Database connection failed: {e}")
            # Continue without database for testing
            self.connection_string = None
            
    async def process_file(self, filepath: str, campaign_name: str = None, client_name: str = None) -> str:
        """
        Process a single file and create semantic chunks
        
        Args:
            filepath: Path to the file to process
            campaign_name: Optional campaign name
            client_name: Optional client name
            
        Returns:
            file_id: UUID of the processed file
        """
        try:
            file_path = Path(filepath)
            if not file_path.exists():
                raise FileNotFoundError(f"File not found: {filepath}")
                
            self.logger.info(f"🔄 Processing file: {file_path.name}")
            
            # Generate file metadata
            file_metadata = self._create_file_metadata(file_path, campaign_name, client_name)
            
            # Extract content based on file type
            chunks = await self._extract_content(file_path, file_metadata)
            
            # Process each chunk with AI
            processed_chunks = []
            for chunk in tqdm(chunks, desc="Processing chunks"):
                processed_chunk = await self._process_chunk(chunk)
                processed_chunks.append(processed_chunk)
                
            # Save to database
            if self.connection_string:
                await self._save_to_database(file_metadata, processed_chunks)
            else:
                # Save to JSON for testing
                await self._save_to_json(file_metadata, processed_chunks)
                
            self.logger.info(f"✅ File processed successfully: {len(processed_chunks)} chunks created")
            return file_metadata.file_id
            
        except Exception as e:
            self.logger.error(f"❌ Error processing file {filepath}: {e}")
            raise
            
    def _create_file_metadata(self, file_path: Path, campaign_name: str = None, client_name: str = None) -> FileMetadata:
        """Create file metadata object"""
        file_id = str(uuid.uuid4())
        file_size = file_path.stat().st_size
        
        # Determine file type and mime type
        suffix = file_path.suffix.lower()
        if suffix in ['.pdf']:
            file_type = 'document'
            mime_type = 'application/pdf'
        elif suffix in ['.pptx', '.ppt']:
            file_type = 'presentation'
            mime_type = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        elif suffix in ['.docx', '.doc']:
            file_type = 'document'
            mime_type = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        elif suffix in ['.jpg', '.jpeg', '.png', '.gif', '.bmp']:
            file_type = 'image'
            mime_type = f'image/{suffix[1:]}'
        elif suffix in ['.mp4', '.avi', '.mov', '.wmv']:
            file_type = 'video'
            mime_type = f'video/{suffix[1:]}'
        else:
            file_type = 'unknown'
            mime_type = 'application/octet-stream'
            
        return FileMetadata(
            file_id=file_id,
            original_filename=file_path.name,
            filepath=str(file_path.absolute()),
            file_size=file_size,
            mime_type=mime_type,
            campaign_name=campaign_name,
            client_name=client_name,
            brand_name=self._extract_brand_name(file_path, campaign_name),
            file_type=file_type,
            processing_status='processing'
        )
        
    def _extract_brand_name(self, file_path: Path, campaign_name: str = None) -> Optional[str]:
        """Extract brand name from file path or campaign name"""
        # Simple brand extraction logic
        common_brands = ['nike', 'adidas', 'coca-cola', 'pepsi', 'mcdonalds', 'kfc', 'jollibee']
        
        text_to_search = f"{file_path.name} {campaign_name or ''}".lower()
        
        for brand in common_brands:
            if brand in text_to_search:
                return brand.title()
                
        return None
        
    async def _extract_content(self, file_path: Path, file_metadata: FileMetadata) -> List[ChunkMetadata]:
        """Extract content from file based on type"""
        chunks = []
        
        if file_metadata.file_type == 'document' and file_path.suffix.lower() == '.pdf':
            chunks = await self._extract_pdf_content(file_path, file_metadata)
        elif file_metadata.file_type == 'presentation':
            chunks = await self._extract_presentation_content(file_path, file_metadata)
        elif file_metadata.file_type == 'image':
            chunks = await self._extract_image_content(file_path, file_metadata)
        elif file_metadata.file_type == 'video':
            chunks = await self._extract_video_content(file_path, file_metadata)
        else:
            # Try to extract as text
            chunks = await self._extract_text_content(file_path, file_metadata)
            
        return chunks
        
    async def _extract_pdf_content(self, file_path: Path, file_metadata: FileMetadata) -> List[ChunkMetadata]:
        """Extract content from PDF file"""
        chunks = []
        
        try:
            doc = fitz.open(str(file_path))
            
            for page_num, page in enumerate(doc):
                # Extract text
                text = page.get_text()
                
                # Extract images
                image_list = page.get_images()
                
                if text.strip():
                    chunk = ChunkMetadata(
                        chunk_id=str(uuid.uuid4()),
                        file_id=file_metadata.file_id,
                        filepath=file_metadata.filepath,
                        title=f"Page {page_num + 1}",
                        text_snippet=text[:1000],  # Limit snippet size
                        tags=[],
                        visual_quality_score=0.0,
                        semantic_topics=[],
                        mood_label='neutral',
                        chunk_index=page_num,
                        chunk_size=len(text),
                        content_type='text',
                        confidence_score=0.8 if len(text) > 50 else 0.4
                    )
                    chunks.append(chunk)
                    
                # Process images on page
                for img_index, img in enumerate(image_list):
                    try:
                        # Extract image
                        xref = img[0]
                        pix = fitz.Pixmap(doc, xref)
                        
                        if pix.n < 5:  # GRAY or RGB
                            img_data = pix.tobytes("png")
                            
                            chunk = ChunkMetadata(
                                chunk_id=str(uuid.uuid4()),
                                file_id=file_metadata.file_id,
                                filepath=file_metadata.filepath,
                                title=f"Page {page_num + 1} - Image {img_index + 1}",
                                text_snippet=f"Image content from page {page_num + 1}",
                                tags=[],
                                visual_quality_score=0.0,
                                semantic_topics=[],
                                mood_label='neutral',
                                chunk_index=page_num * 100 + img_index,  # Unique index
                                chunk_size=len(img_data),
                                content_type='image',
                                confidence_score=0.7
                            )
                            chunks.append(chunk)
                            
                        pix = None
                    except Exception as e:
                        self.logger.warning(f"Error extracting image from page {page_num}: {e}")
                        
            doc.close()
            
        except Exception as e:
            self.logger.error(f"Error extracting PDF content: {e}")
            
        return chunks
        
    async def _extract_presentation_content(self, file_path: Path, file_metadata: FileMetadata) -> List[ChunkMetadata]:
        """Extract content from PowerPoint presentation"""
        chunks = []
        
        try:
            prs = Presentation(str(file_path))
            
            for slide_index, slide in enumerate(prs.slides):
                # Extract text from slide
                text_runs = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
                        
                slide_text = " ".join(text_runs)
                
                if slide_text.strip():
                    chunk = ChunkMetadata(
                        chunk_id=str(uuid.uuid4()),
                        file_id=file_metadata.file_id,
                        filepath=file_metadata.filepath,
                        title=f"Slide {slide_index + 1}",
                        text_snippet=slide_text[:1000],
                        tags=[],
                        visual_quality_score=0.0,
                        semantic_topics=[],
                        mood_label='neutral',
                        chunk_index=slide_index,
                        chunk_size=len(slide_text),
                        content_type='slide',
                        confidence_score=0.8 if len(slide_text) > 20 else 0.5
                    )
                    chunks.append(chunk)
                    
        except Exception as e:
            self.logger.error(f"Error extracting presentation content: {e}")
            
        return chunks
        
    async def _extract_image_content(self, file_path: Path, file_metadata: FileMetadata) -> List[ChunkMetadata]:
        """Extract content from image file"""
        chunks = []
        
        try:
            # Load image
            img = Image.open(str(file_path))
            
            # Basic image analysis
            width, height = img.size
            mode = img.mode
            
            chunk = ChunkMetadata(
                chunk_id=str(uuid.uuid4()),
                file_id=file_metadata.file_id,
                filepath=file_metadata.filepath,
                title=file_path.stem,
                text_snippet=f"Image: {width}x{height} pixels, {mode} mode",
                tags=[],
                visual_quality_score=0.0,  # Will be calculated by AI
                semantic_topics=[],
                mood_label='neutral',
                chunk_index=0,
                chunk_size=file_metadata.file_size,
                content_type='image',
                confidence_score=0.9
            )
            chunks.append(chunk)
            
        except Exception as e:
            self.logger.error(f"Error extracting image content: {e}")
            
        return chunks
        
    async def _extract_video_content(self, file_path: Path, file_metadata: FileMetadata) -> List[ChunkMetadata]:
        """Extract content from video file (extract frames)"""
        chunks = []
        
        try:
            # Extract key frames using OpenCV
            cap = cv2.VideoCapture(str(file_path))
            
            if not cap.isOpened():
                raise ValueError(f"Cannot open video file: {file_path}")
                
            fps = cap.get(cv2.CAP_PROP_FPS)
            frame_count = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
            duration = frame_count / fps if fps > 0 else 0
            
            # Extract frames every 10 seconds
            frames_to_extract = max(1, int(duration / 10))
            frame_interval = max(1, frame_count // frames_to_extract)
            
            frame_index = 0
            extracted_frames = 0
            
            while cap.isOpened() and extracted_frames < 10:  # Limit to 10 frames
                ret, frame = cap.read()
                if not ret:
                    break
                    
                if frame_index % frame_interval == 0:
                    timestamp = frame_index / fps if fps > 0 else 0
                    
                    chunk = ChunkMetadata(
                        chunk_id=str(uuid.uuid4()),
                        file_id=file_metadata.file_id,
                        filepath=file_metadata.filepath,
                        title=f"Frame at {timestamp:.1f}s",
                        text_snippet=f"Video frame at {timestamp:.1f} seconds",
                        tags=[],
                        visual_quality_score=0.0,
                        semantic_topics=[],
                        mood_label='neutral',
                        chunk_index=extracted_frames,
                        chunk_size=frame.nbytes,
                        content_type='video',
                        confidence_score=0.7
                    )
                    chunks.append(chunk)
                    extracted_frames += 1
                    
                frame_index += 1
                
            cap.release()
            
        except Exception as e:
            self.logger.error(f"Error extracting video content: {e}")
            
        return chunks
        
    async def _extract_text_content(self, file_path: Path, file_metadata: FileMetadata) -> List[ChunkMetadata]:
        """Extract content from text file"""
        chunks = []
        
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
                
            # Split into chunks of ~1000 characters
            chunk_size = 1000
            for i in range(0, len(content), chunk_size):
                chunk_text = content[i:i + chunk_size]
                
                chunk = ChunkMetadata(
                    chunk_id=str(uuid.uuid4()),
                    file_id=file_metadata.file_id,
                    filepath=file_metadata.filepath,
                    title=f"Text chunk {i // chunk_size + 1}",
                    text_snippet=chunk_text[:500],
                    tags=[],
                    visual_quality_score=0.0,
                    semantic_topics=[],
                    mood_label='neutral',
                    chunk_index=i // chunk_size,
                    chunk_size=len(chunk_text),
                    content_type='text',
                    confidence_score=0.8
                )
                chunks.append(chunk)
                
        except Exception as e:
            self.logger.error(f"Error extracting text content: {e}")
            
        return chunks
        
    async def _process_chunk(self, chunk: ChunkMetadata) -> ChunkMetadata:
        """Process chunk with AI analysis"""
        try:
            # Generate embedding
            embedding = await self._generate_embedding(chunk.text_snippet)
            chunk.embedding_vector = embedding
            
            # Classify mood and topics
            classification = await self._classify_content(chunk.text_snippet, chunk.content_type)
            chunk.mood_label = classification.get('mood', 'neutral')
            chunk.semantic_topics = classification.get('topics', [])
            chunk.tags = classification.get('tags', [])
            
            # Calculate visual quality score (for images/videos)
            if chunk.content_type in ['image', 'video', 'slide']:
                chunk.visual_quality_score = await self._assess_visual_quality(chunk)
            else:
                chunk.visual_quality_score = 0.5  # Default for text
                
        except Exception as e:
            self.logger.warning(f"Error processing chunk {chunk.chunk_id}: {e}")
            # Set defaults
            chunk.embedding_vector = [0.0] * 1536  # Default embedding size
            chunk.mood_label = 'neutral'
            chunk.semantic_topics = []
            chunk.tags = []
            chunk.visual_quality_score = 0.5
            
        return chunk
        
    async def _generate_embedding(self, text: str) -> List[float]:
        """Generate embedding vector for text"""
        try:
            if not self.openai_client:
                return [0.0] * 1536  # Default embedding size
                
            response = self.openai_client.embeddings.create(
                input=text,
                model=os.getenv('AZURE_OPENAI_EMBEDDING_MODEL', 'text-embedding-ada-002')
            )
            
            return response.data[0].embedding
            
        except Exception as e:
            self.logger.warning(f"Error generating embedding: {e}")
            return [0.0] * 1536
            
    async def _classify_content(self, text: str, content_type: str) -> Dict[str, Any]:
        """Classify content mood, topics, and tags"""
        try:
            if not self.openai_client:
                return {'mood': 'neutral', 'topics': [], 'tags': []}
                
            prompt = f"""
            Analyze this {content_type} content and provide:
            1. Mood (happy, sad, neutral, energetic, calm, professional, creative, urgent)
            2. Topics (max 5 key topics)
            3. Tags (max 10 relevant tags)
            
            Content: {text[:500]}
            
            Respond in JSON format:
            {{
                "mood": "mood_label",
                "topics": ["topic1", "topic2"],
                "tags": ["tag1", "tag2"]
            }}
            """
            
            response = self.openai_client.chat.completions.create(
                model=os.getenv('AZURE_OPENAI_CHAT_MODEL', 'gpt-4'),
                messages=[{"role": "user", "content": prompt}],
                max_tokens=200,
                temperature=0.3
            )
            
            result = json.loads(response.choices[0].message.content)
            return result
            
        except Exception as e:
            self.logger.warning(f"Error classifying content: {e}")
            return {'mood': 'neutral', 'topics': [], 'tags': []}
            
    async def _assess_visual_quality(self, chunk: ChunkMetadata) -> float:
        """Assess visual quality of image/video content"""
        try:
            # Simple quality assessment
            # In a real implementation, this would use computer vision models
            
            if chunk.content_type == 'image':
                # For images, use file size and content type as proxy
                if chunk.chunk_size > 100000:  # > 100KB
                    return 0.8
                elif chunk.chunk_size > 50000:  # > 50KB
                    return 0.6
                else:
                    return 0.4
            elif chunk.content_type in ['video', 'slide']:
                return 0.7  # Default for video/slides
            else:
                return 0.5  # Default
                
        except Exception as e:
            self.logger.warning(f"Error assessing visual quality: {e}")
            return 0.5
            
    async def _save_to_database(self, file_metadata: FileMetadata, chunks: List[ChunkMetadata]):
        """Save processed data to SQL Server database"""
        try:
            with pyodbc.connect(self.connection_string) as conn:
                cursor = conn.cursor()
                
                # Insert file metadata
                cursor.execute("""
                    INSERT INTO fileMetadata (
                        file_id, original_filename, filepath, file_size, mime_type,
                        campaign_name, client_name, brand_name, file_type, 
                        processing_status, total_chunks, processed_at
                    ) VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
                """, (
                    file_metadata.file_id, file_metadata.original_filename,
                    file_metadata.filepath, file_metadata.file_size, file_metadata.mime_type,
                    file_metadata.campaign_name, file_metadata.client_name, file_metadata.brand_name,
                    file_metadata.file_type, 'completed', len(chunks), datetime.utcnow()
                ))
                
                # Insert chunks
                for chunk in chunks:
                    cursor.execute("""
                        INSERT INTO pageIndex (
                            chunk_id, file_id, filepath, title, text_snippet, tags,
                            visual_quality_score, semantic_topics, mood_label,
                            embedding_vector, chunk_index, chunk_size, content_type, confidence_score
                        ) VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
                    """, (
                        chunk.chunk_id, chunk.file_id, chunk.filepath, chunk.title,
                        chunk.text_snippet, json.dumps(chunk.tags), chunk.visual_quality_score,
                        json.dumps(chunk.semantic_topics), chunk.mood_label,
                        json.dumps(chunk.embedding_vector) if chunk.embedding_vector else None,
                        chunk.chunk_index, chunk.chunk_size, chunk.content_type, chunk.confidence_score
                    ))
                
                conn.commit()
                self.logger.info(f"✅ Saved to database: {len(chunks)} chunks")
                
        except Exception as e:
            self.logger.error(f"❌ Database save failed: {e}")
            # Fallback to JSON
            await self._save_to_json(file_metadata, chunks)
            
    async def _save_to_json(self, file_metadata: FileMetadata, chunks: List[ChunkMetadata]):
        """Save processed data to JSON file (fallback)"""
        try:
            output_dir = Path("pageindex_output")
            output_dir.mkdir(exist_ok=True)
            
            output_file = output_dir / f"{file_metadata.file_id}.json"
            
            data = {
                "file_metadata": asdict(file_metadata),
                "chunks": [asdict(chunk) for chunk in chunks],
                "processed_at": datetime.utcnow().isoformat(),
                "agent_version": "1.0.0"
            }
            
            with open(output_file, 'w', encoding='utf-8') as f:
                json.dump(data, f, indent=2, ensure_ascii=False)
                
            self.logger.info(f"✅ Saved to JSON: {output_file}")
            
        except Exception as e:
            self.logger.error(f"❌ JSON save failed: {e}")

async def main():
    """CLI entry point"""
    parser = argparse.ArgumentParser(description='TBWA PageIndex Agent - Semantic file processing')
    parser.add_argument('filepath', help='Path to file to process')
    parser.add_argument('--campaign', help='Campaign name')
    parser.add_argument('--client', help='Client name')
    parser.add_argument('--batch', action='store_true', help='Process all files in directory')
    parser.add_argument('--debug', action='store_true', help='Enable debug logging')
    
    args = parser.parse_args()
    
    if args.debug:
        logging.getLogger().setLevel(logging.DEBUG)
        
    agent = PageIndexAgent()
    
    try:
        if args.batch:
            # Process all files in directory
            file_path = Path(args.filepath)
            if not file_path.is_dir():
                raise ValueError("Batch mode requires a directory path")
                
            files = list(file_path.rglob('*'))
            files = [f for f in files if f.is_file() and f.suffix.lower() in 
                    ['.pdf', '.pptx', '.docx', '.jpg', '.jpeg', '.png', '.mp4', '.mov']]
                    
            print(f"🔄 Processing {len(files)} files in batch mode...")
            
            for file in tqdm(files, desc="Processing files"):
                try:
                    await agent.process_file(str(file), args.campaign, args.client)
                except Exception as e:
                    print(f"❌ Error processing {file}: {e}")
                    
        else:
            # Process single file
            file_id = await agent.process_file(args.filepath, args.campaign, args.client)
            print(f"✅ File processed successfully. File ID: {file_id}")
            
    except Exception as e:
        print(f"❌ Processing failed: {e}")
        sys.exit(1)

if __name__ == "__main__":
    asyncio.run(main())